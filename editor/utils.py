import os
import pypdf
from pptx import Presentation
from docx import Document as DocxDocument
from docx.shared import Inches
import mammoth
from django.core.files.base import ContentFile
from django.conf import settings

def get_file_type(filename):
    ext = os.path.splitext(filename)[1].lower()
    if ext == '.csv':
        return 'csv'
    elif ext == '.pdf':
        return 'pdf'
    elif ext in ['.docx', '.doc']:
        return 'docx'
    elif ext in ['.pptx', '.ppt']:
        return 'pptx'
    return 'other'

def merge_pdfs(file_paths, output_path):
    merger = pypdf.PdfWriter()
    for path in file_paths:
        merger.append(path)
    merger.write(output_path)
    merger.close()
    return output_path

def split_pdf(file_path, output_dir, pages=None):
    reader = pypdf.PdfReader(file_path)
    output_files = []
    
    # If pages not specified, split each page
    if not pages:
        for i in range(len(reader.pages)):
            writer = pypdf.PdfWriter()
            writer.add_page(reader.pages[i])
            out_filename = f"page_{i+1}.pdf"
            out_path = os.path.join(output_dir, out_filename)
            with open(out_path, "wb") as f:
                writer.write(f)
            output_files.append(out_path)
    return output_files

def replace_text_docx(file_path, search_text, replace_text, output_path):
    doc = DocxDocument(file_path)
    for paragraph in doc.paragraphs:
        if search_text in paragraph.text:
            paragraph.text = paragraph.text.replace(search_text, replace_text)
    
    # Also check tables
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    if search_text in paragraph.text:
                        paragraph.text = paragraph.text.replace(search_text, replace_text)
                        
    doc.save(output_path)
    return output_path

def replace_text_pptx(file_path, search_text, replace_text, output_path):
    prs = Presentation(file_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                if search_text in shape.text:
                    shape.text = shape.text.replace(search_text, replace_text)
    prs.save(output_path)
    prs.save(output_path)
    return output_path

def organize_pdf(file_path, page_config, output_path):
    """
    Reorder, rotate, and select pages from a PDF.
    page_config: List of dicts [{'page': 0, 'rotate': 90}, ...]
    Note: page_config uses 0-based indexing for internal logic, but frontend might send 1-based.
    """
    reader = pypdf.PdfReader(file_path)
    writer = pypdf.PdfWriter()
    
    for config in page_config:
        page_idx = config.get('original_page_number') - 1 # Convert 1-based to 0-based
        rotation = config.get('rotate', 0)
        
        if 0 <= page_idx < len(reader.pages):
            page = reader.pages[page_idx]
            if rotation:
                page.rotate(rotation)
            writer.add_page(page)
            
    with open(output_path, "wb") as f:
        writer.write(f)
    return output_path

import mammoth

def docx_to_html(file_path):
    """Convert DOCX to HTML for preview"""
    with open(file_path, "rb") as docx_file:
        result = mammoth.convert_to_html(docx_file)
        return result.value

def get_docx_paragraphs(file_path):
    """Get all paragraphs from DOCX for editing"""
    doc = DocxDocument(file_path)
    paragraphs = []
    for i, p in enumerate(doc.paragraphs):
        if p.text.strip(): # Only non-empty
            paragraphs.append({'id': i, 'text': p.text})
    return paragraphs

def update_docx_paragraph(file_path, index, new_text):
    """Update a specific paragraph"""
    doc = DocxDocument(file_path)
    if 0 <= index < len(doc.paragraphs):
        doc.paragraphs[index].text = new_text
        doc.save(file_path)
        return True
    return False

def add_image_to_docx(file_path, image_file):
    """Append image to end of DOCX"""
    doc = DocxDocument(file_path)
    doc.add_picture(image_file, width=None) # Auto width or set inches
    doc.save(file_path)
    return True

from docxcompose.composer import Composer

def merge_docx_files(file_paths, output_path):
    """Merge multiple DOCX files into one"""
    if not file_paths:
        return None
        
    master = DocxDocument(file_paths[0])
    composer = Composer(master)
    
    for i in range(1, len(file_paths)):
        doc2 = DocxDocument(file_paths[i])
        composer.append(doc2)
        
    composer.save(output_path)
    return output_path

def organize_docx_pages(file_path, page_config, output_path):
    """
    Split/Reorder/Delete 'pages' based on hard page breaks.
    page_config: List of dicts [{'original_page_number': 1}, ...]
    Strategy: For each target page, load original, delete others, save temp. Then merge.
    """
    import copy
    import tempfile
    
    # 1. Analyze original to find page boundaries (element indices)
    doc = DocxDocument(file_path)
    body_elements = list(doc.element.body)
    
    page_boundaries = [] # List of (start_index, end_index)
    current_start = 0
    
    for i, element in enumerate(body_elements):
        # Check for page break
        has_break = False
        if element.tag.endswith('p'):
            if 'w:type="page"' in element.xml:
                has_break = True
        
        if has_break:
            # The break is usually AT THE END of the page content, or creates a new one.
            # A hard break ending a page means this paragraph is the last one of the current page.
            page_boundaries.append((current_start, i)) # Inclusive start, Inclusive end
            current_start = i + 1
            
    # Add last page
    if current_start < len(body_elements):
        page_boundaries.append((current_start, len(body_elements) - 1))
        
    num_original_pages = len(page_boundaries)
    temp_files = []
    
    try:
        # 2. Extract specific pages
        for config in page_config:
            page_idx = config.get('original_page_number') - 1
            
            if 0 <= page_idx < num_original_pages:
                # Load fresh copy to preserve styles
                temp_doc = DocxDocument(file_path)
                
                # We need to remove elements that are NOT in the desired range.
                # It's safer to remove from end to start to avoid index shifting, 
                # but since we are removing based on original indices, we must be careful.
                # Actually, the indices in temp_doc.element.body should match body_elements initially.
                
                # Get range to KEEP
                start, end = page_boundaries[page_idx]
                
                # Identify elements to DELETE
                # Delete everything after 'end'
                # Delete everything before 'start'
                
                # We can't use slice assignment on xml body easily.
                # Iterate backwards and remove if index not in range [start, end]
                
                # Note: body_elements includes section properties etc? 
                # doc.element.body children include paragraphs, tables, sectPr.
                # We should keep sectPr (usually last) to preserve margins? 
                # If we delete sectPr, we might lose layout. 
                # For now, let's just delete P and TBL.
                
                current_body = list(temp_doc.element.body)
                for i in range(len(current_body) - 1, -1, -1):
                    # Keep the last element if it is sectPr?
                    # valid content range is [start, end].
                    if i < start or i > end:
                         # Don't delete sectPr if it's the very last element and we are just removing content
                         # But a page break might define a section. 
                         # Let's simplified: just remove if not in range.
                         # If we remove sectPr, it defaults.
                         
                         el = current_body[i]
                         if not el.tag.endswith('sectPr'):
                             el.getparent().remove(el)
                
                # Save temp chunk
                with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as tmp:
                    temp_doc.save(tmp.name)
                    temp_files.append(tmp.name)
        
        # 3. Merge chunks
        if not temp_files:
            return None
            
        merge_docx_files(temp_files, output_path)
        return output_path
        
    finally:
        # Cleanup
        for path in temp_files:
            try:
                os.remove(path)
            except: pass 

