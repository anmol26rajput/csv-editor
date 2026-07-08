"""PDF conversion service: turn an uploaded PDF into other document/image formats.

Every converter takes the source path plus an output directory and returns
(output_path, file_type). Multi-file outputs (per-page images/SVGs) are zipped.
"""
import html as html_lib
import io
import os
import shutil
import subprocess
import zipfile

import fitz  # PyMuPDF

# Rendering resolution for raster outputs (images, PPTX slides)
RASTER_DPI = 150

SUPPORTED_FORMATS = [
    'docx', 'word', 'xlsx', 'excel', 'pptx',
    'txt', 'text', 'html', 'svg', 'epub',
    'png', 'jpg', 'webp', 'tiff', 'image',
]

# Source types accepted by convert_to_pdf
IMAGE_SOURCE_EXTS = {'png', 'jpg', 'jpeg', 'webp', 'tiff', 'tif', 'bmp', 'gif'}
OFFICE_SOURCE_EXTS = {'docx', 'doc', 'xlsx', 'xls', 'pptx', 'ppt', 'odt', 'ods', 'odp', 'rtf', 'html', 'htm', 'csv'}
TEXT_SOURCE_EXTS = {'txt', 'md'}
TO_PDF_SOURCE_EXTS = IMAGE_SOURCE_EXTS | OFFICE_SOURCE_EXTS | TEXT_SOURCE_EXTS


def convert_pdf(input_path, output_dir, target):
    base_name = os.path.splitext(os.path.basename(input_path))[0]
    target = target.lower()

    converters = {
        'docx': _to_docx, 'word': _to_docx,
        'xlsx': _to_xlsx, 'excel': _to_xlsx,
        'pptx': _to_pptx,
        'txt': _to_txt, 'text': _to_txt,
        'html': _to_html,
        'svg': _to_svg,
        'epub': _to_epub,
        'png': lambda i, o, b: _to_images(i, o, b, 'png'),
        'jpg': lambda i, o, b: _to_images(i, o, b, 'jpg'),
        'webp': lambda i, o, b: _to_images(i, o, b, 'webp'),
        'image': lambda i, o, b: _to_images(i, o, b, 'png'),
        'tiff': _to_tiff,
    }
    if target not in converters:
        raise ValueError(f"Unsupported target format: {target}")
    return converters[target](input_path, output_dir, base_name)


def _to_docx(input_path, output_dir, base_name):
    from pdf2docx import Converter
    output_path = os.path.join(output_dir, f"{base_name}.docx")
    cv = Converter(input_path)
    try:
        cv.convert(output_path)
    finally:
        cv.close()
    return output_path, 'docx'


def _to_xlsx(input_path, output_dir, base_name):
    from openpyxl import Workbook
    output_path = os.path.join(output_dir, f"{base_name}.xlsx")
    wb = Workbook()
    wb.remove(wb.active)

    with fitz.open(input_path) as doc:
        for page_index, page in enumerate(doc):
            ws = wb.create_sheet(title=f"Page {page_index + 1}")
            tables = page.find_tables()
            row_cursor = 1
            if tables.tables:
                for table in tables:
                    for row in table.extract():
                        for col_index, cell in enumerate(row, start=1):
                            ws.cell(row=row_cursor, column=col_index, value=cell)
                        row_cursor += 1
                    row_cursor += 1  # blank row between tables
            else:
                # No detectable tables: fall back to one text line per row
                for line in page.get_text().splitlines():
                    if line.strip():
                        ws.cell(row=row_cursor, column=1, value=line)
                        row_cursor += 1
    wb.save(output_path)
    return output_path, 'xlsx'


def _to_pptx(input_path, output_dir, base_name):
    from pptx import Presentation
    from pptx.util import Emu

    output_path = os.path.join(output_dir, f"{base_name}.pptx")
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]

    with fitz.open(input_path) as doc:
        first = doc[0].rect
        # Match slide size to the PDF page aspect ratio (1pt = 12700 EMU)
        prs.slide_width = Emu(int(first.width * 12700))
        prs.slide_height = Emu(int(first.height * 12700))

        for page in doc:
            pix = page.get_pixmap(dpi=RASTER_DPI)
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(
                io.BytesIO(pix.tobytes('png')), 0, 0,
                width=prs.slide_width, height=prs.slide_height
            )
    prs.save(output_path)
    return output_path, 'pptx'


def _to_txt(input_path, output_dir, base_name):
    output_path = os.path.join(output_dir, f"{base_name}.txt")
    with fitz.open(input_path) as doc, open(output_path, 'w', encoding='utf-8') as f:
        for page_index, page in enumerate(doc):
            if page_index:
                f.write('\n\n')
            f.write(page.get_text())
    return output_path, 'txt'


def _to_html(input_path, output_dir, base_name):
    output_path = os.path.join(output_dir, f"{base_name}.html")
    with fitz.open(input_path) as doc:
        body = '\n'.join(page.get_text('xhtml') for page in doc)
    html = (
        '<!DOCTYPE html>\n<html>\n<head>\n<meta charset="utf-8">\n'
        f'<title>{base_name}</title>\n</head>\n<body>\n{body}\n</body>\n</html>\n'
    )
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(html)
    return output_path, 'html'


def _to_svg(input_path, output_dir, base_name):
    with fitz.open(input_path) as doc:
        svgs = [page.get_svg_image() for page in doc]

    if len(svgs) == 1:
        output_path = os.path.join(output_dir, f"{base_name}.svg")
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(svgs[0])
        return output_path, 'svg'

    output_path = os.path.join(output_dir, f"{base_name}_svg.zip")
    with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zf:
        for i, svg in enumerate(svgs):
            zf.writestr(f"{base_name}_page_{i + 1}.svg", svg)
    return output_path, 'zip'


def _to_epub(input_path, output_dir, base_name):
    from ebooklib import epub

    output_path = os.path.join(output_dir, f"{base_name}.epub")
    book = epub.EpubBook()
    book.set_identifier(base_name)
    book.set_title(base_name)
    book.set_language('en')

    chapters = []
    with fitz.open(input_path) as doc:
        for page_index, page in enumerate(doc):
            chapter = epub.EpubHtml(
                title=f"Page {page_index + 1}",
                file_name=f"page_{page_index + 1}.xhtml",
                lang='en'
            )
            chapter.content = f"<html><body>{page.get_text('xhtml')}</body></html>"
            book.add_item(chapter)
            chapters.append(chapter)

    book.toc = chapters
    book.spine = ['nav'] + chapters
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())
    epub.write_epub(output_path, book)
    return output_path, 'epub'


def _to_images(input_path, output_dir, base_name, ext):
    from PIL import Image

    def save_page(page, path):
        pix = page.get_pixmap(dpi=RASTER_DPI)
        if ext == 'webp':
            Image.open(io.BytesIO(pix.tobytes('png'))).save(path, 'WEBP')
        else:
            pix.save(path)

    with fitz.open(input_path) as doc:
        if doc.page_count == 1:
            output_path = os.path.join(output_dir, f"{base_name}.{ext}")
            save_page(doc[0], output_path)
            return output_path, ext

        output_path = os.path.join(output_dir, f"{base_name}_{ext}.zip")
        with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zf:
            for i, page in enumerate(doc):
                page_path = os.path.join(output_dir, f"{base_name}_page_{i + 1}.{ext}")
                save_page(page, page_path)
                zf.write(page_path, os.path.basename(page_path))
                os.remove(page_path)
        return output_path, 'zip'


def _to_tiff(input_path, output_dir, base_name):
    from PIL import Image

    output_path = os.path.join(output_dir, f"{base_name}.tiff")
    with fitz.open(input_path) as doc:
        frames = [
            Image.open(io.BytesIO(page.get_pixmap(dpi=RASTER_DPI).tobytes('png')))
            for page in doc
        ]
    frames[0].save(output_path, 'TIFF', save_all=True, append_images=frames[1:], compression='tiff_deflate')
    return output_path, 'tiff'


def convert_to_pdf(input_path, output_dir):
    """Convert a non-PDF source file into a PDF. Returns the output path."""
    base_name, ext = os.path.splitext(os.path.basename(input_path))
    ext = ext[1:].lower()
    output_path = os.path.join(output_dir, f"{base_name}.pdf")

    if ext in IMAGE_SOURCE_EXTS:
        return _image_to_pdf(input_path, output_path)
    if ext in TEXT_SOURCE_EXTS:
        return _text_to_pdf(input_path, output_path)
    if ext in OFFICE_SOURCE_EXTS:
        return _office_to_pdf(input_path, output_dir, base_name)
    raise ValueError(f"Cannot convert '.{ext}' files to PDF")


def _image_to_pdf(input_path, output_path):
    from PIL import Image, ImageSequence

    with Image.open(input_path) as img:
        frames = []
        for frame in ImageSequence.Iterator(img):
            frame = frame.convert('RGB')
            frames.append(frame.copy())
    frames[0].save(output_path, 'PDF', save_all=True, append_images=frames[1:])
    return output_path


def _text_to_pdf(input_path, output_path):
    with open(input_path, 'r', encoding='utf-8', errors='replace') as f:
        text = f.read()

    paragraphs = ''.join(
        f"<p>{html_lib.escape(line) or '&nbsp;'}</p>"
        for line in text.splitlines()
    )
    story_html = (
        '<html><body style="font-family: sans-serif; font-size: 11px;">'
        f'{paragraphs}</body></html>'
    )

    story = fitz.Story(html=story_html)
    writer = fitz.DocumentWriter(output_path)
    mediabox = fitz.paper_rect('a4')
    content_box = mediabox + (36, 36, -36, -36)  # 0.5 inch margins

    more = True
    while more:
        device = writer.begin_page(mediabox)
        more, _ = story.place(content_box)
        story.draw(device)
        writer.end_page()
    writer.close()
    return output_path


def _office_to_pdf(input_path, output_dir, base_name):
    if not shutil.which('soffice'):
        raise ValueError(
            'Converting Office documents to PDF requires LibreOffice, '
            'which is not installed on the server.'
        )

    cmd = [
        'soffice', '--headless',
        '--convert-to', 'pdf',
        '--outdir', output_dir,
        input_path,
    ]
    result = subprocess.run(cmd, capture_output=True, timeout=120)
    output_path = os.path.join(output_dir, f"{base_name}.pdf")
    if result.returncode != 0 or not os.path.exists(output_path):
        stderr = result.stderr.decode('utf-8', errors='replace').strip()
        raise ValueError(f"LibreOffice conversion failed: {stderr or 'no output produced'}")
    return output_path
